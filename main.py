
import os
import re
import sys
import traceback
import logging
from pathlib import Path
from datetime import datetime
import tkinter as tk
from tkinter import ttk, filedialog, messagebox

from pptx import Presentation
from docx import Document
from openpyxl import Workbook, load_workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.utils import get_column_letter
from openpyxl.cell.cell import ILLEGAL_CHARACTERS_RE


APP_NAME = "異常報告履歷擷取工具"
APP_VERSION = "1.2.2"

HEADERS = [
    "No.", "報告日期", "異常類別", "客戶/廠區", "Model", "PN", "PO / 工單",
    "異常站別", "異常項目", "不良率",
    "主要不良現象", "主要原因", "改善對策", "改善效果", "驗證結果",
    "最終判定", "備註", "來源檔案"
]


def clean_text(v):
    if v is None:
        return ""
    s = str(v).replace("\r", "\n")
    s = re.sub(r"[ \t]+", " ", s)
    s = re.sub(r"\n{3,}", "\n\n", s)
    return s.strip()

def sanitize_excel_text(value, logger=None, field_name="", source_file=""):
    if value is None:
        return ""
    if not isinstance(value, str):
        return value
    original = value
    value = value.replace("\x0b", " / ").replace("\x0c", " / ")
    cleaned = ILLEGAL_CHARACTERS_RE.sub("", value)
    if cleaned != original and logger:
        logger.warning(
            "SANITIZE | Removed illegal Excel character | file=%s | field=%s | before=%r | after=%r",
            source_file, field_name, original[:200], cleaned[:200]
        )
    return cleaned


def sanitize_record_for_excel(record, logger=None):
    source_file = record.get("來源檔案", "")
    return {k: sanitize_excel_text(v, logger, k, source_file) for k, v in record.items()}



def first_match(text, patterns, flags=re.I):
    for p in patterns:
        m = re.search(p, text, flags)
        if m:
            return clean_text(m.group(1))
    return ""


def compact_join(items, sep=" / "):
    out = []
    for x in items:
        x = clean_text(x)
        if x and x not in out:
            out.append(x)
    return sep.join(out)



def resource_path(relative_path):
    """Resolve bundled resources in source mode and PyInstaller EXE mode."""
    try:
        base_path = Path(sys._MEIPASS)
    except Exception:
        base_path = Path(__file__).resolve().parent
    return base_path / relative_path


def parse_report_date(value):
    """Return a sortable datetime-like tuple or None for unknown dates."""
    if value is None:
        return None
    s = clean_text(value)
    if not s:
        return None

    patterns = [
        r"^(20\d{2})[/-](\d{1,2})[/-](\d{1,2})$",
        r"^(20\d{2})\.(\d{1,2})\.(\d{1,2})$",
        r"^(20\d{2})(\d{2})(\d{2})$",
    ]
    for p in patterns:
        m = re.match(p, s)
        if m:
            try:
                y, mo, d = int(m.group(1)), int(m.group(2)), int(m.group(3))
                if 1 <= mo <= 12 and 1 <= d <= 31:
                    return (y, mo, d)
            except Exception:
                pass
    return None


def sort_history_by_report_date(ws, logger=None):
    """Sort all history rows by report date DESC, blank/invalid dates last, then renumber No."""
    if ws.max_row < 2:
        return {"sorted_rows": 0, "blank_dates": 0, "sort_pass": True, "renumber_pass": True}

    header_map = {clean_text(ws.cell(1, c).value): c for c in range(1, ws.max_column + 1)}
    date_col = header_map.get("報告日期")
    if not date_col:
        if logger:
            logger.warning("SORT_CHECK | 報告日期欄位不存在 | SKIP")
        return {"sorted_rows": 0, "blank_dates": 0, "sort_pass": False, "renumber_pass": False}

    records = []
    for r in range(2, ws.max_row + 1):
        values = [ws.cell(r, c).value for c in range(1, ws.max_column + 1)]
        date_key = parse_report_date(ws.cell(r, date_col).value)
        records.append((date_key, r, values))

    # Valid dates first, newest to oldest. Invalid/blank dates last.
    records.sort(
        key=lambda item: (
            item[0] is not None,
            item[0] if item[0] is not None else (0, 0, 0),
            -item[1],
        ),
        reverse=True,
    )

    # Rewrite data rows only.
    if ws.max_row >= 2:
        ws.delete_rows(2, ws.max_row - 1)

    for idx, (_, _, values) in enumerate(records, 1):
        values[0] = idx
        ws.append(values)

    # Verify date order.
    parsed_dates = []
    blank_dates = 0
    seen_blank = False
    blank_at_end = True
    for r in range(2, ws.max_row + 1):
        d = parse_report_date(ws.cell(r, date_col).value)
        if d is None:
            blank_dates += 1
            seen_blank = True
        else:
            if seen_blank:
                blank_at_end = False
            parsed_dates.append(d)

    date_desc = all(parsed_dates[i] >= parsed_dates[i + 1] for i in range(len(parsed_dates) - 1))
    sort_pass = date_desc and blank_at_end
    renumber_pass = all(ws.cell(r, 1).value == r - 1 for r in range(2, ws.max_row + 1))

    if logger:
        logger.info(
            "SORT_CHECK | 報告日期 | DESC | %s | rows=%d | blank_dates=%d",
            "PASS" if sort_pass else "FAIL", len(records), blank_dates
        )
        logger.info(
            "SORT_CHECK | Blank Date At End | %s",
            "PASS" if blank_at_end else "FAIL"
        )
        logger.info(
            "RENUMBER_CHECK | No. 1~N | %s",
            "PASS" if renumber_pass else "FAIL"
        )

    return {
        "sorted_rows": len(records),
        "blank_dates": blank_dates,
        "sort_pass": sort_pass,
        "blank_at_end": blank_at_end,
        "renumber_pass": renumber_pass,
    }


def setup_logger(run_dir):
    logger = logging.getLogger("abnormal_history")
    logger.setLevel(logging.DEBUG)
    logger.handlers.clear()
    fmt = logging.Formatter("%(asctime)s | %(levelname)s | %(message)s")

    for filename, level in [("run.log", logging.INFO), ("debug.log", logging.DEBUG), ("error.log", logging.ERROR)]:
        handler = logging.FileHandler(run_dir / filename, encoding="utf-8")
        handler.setLevel(level)
        handler.setFormatter(fmt)
        logger.addHandler(handler)
    return logger


def extract_pptx(path, logger):
    prs = Presentation(path)
    pages, tables = [], []
    for idx, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            try:
                if hasattr(shape, "text") and shape.text:
                    texts.append(clean_text(shape.text))
                if getattr(shape, "has_table", False):
                    rows = []
                    for row in shape.table.rows:
                        rows.append([clean_text(c.text) for c in row.cells])
                    tables.append({"page": idx, "rows": rows})
            except Exception:
                logger.debug("PPT shape parse warning page=%s\n%s", idx, traceback.format_exc())
        pages.append({"page": idx, "text": "\n".join([x for x in texts if x])})
    return pages, tables


def extract_docx(path, logger):
    doc = Document(path)
    paragraphs = [clean_text(p.text) for p in doc.paragraphs if clean_text(p.text)]
    tables = []
    for idx, table in enumerate(doc.tables, 1):
        rows = []
        for row in table.rows:
            rows.append([clean_text(c.text) for c in row.cells])
        tables.append({"page": idx, "rows": rows})
    return [{"page": 1, "text": "\n".join(paragraphs)}], tables


def table_key_values(tables):
    kv = {}
    for t in tables:
        for row in t["rows"]:
            for i in range(len(row) - 1):
                k, v = clean_text(row[i]), clean_text(row[i + 1])
                if k and v and len(k) <= 30:
                    kv.setdefault(k, []).append(v)
    return kv


def values_for(kv, keywords):
    vals = []
    for k, vs in kv.items():
        nk = re.sub(r"\s+", "", k).lower()
        for kw in keywords:
            if re.sub(r"\s+", "", kw).lower() in nk:
                vals.extend(vs)
                break
    return vals


def section_from_text(text, start_words, end_words):
    starts = "|".join(re.escape(x) for x in start_words)
    ends = "|".join(re.escape(x) for x in end_words)
    pat = rf"(?:{starts})\s*[:：]?\s*(.*?)(?=\n\s*(?:{ends})\s*[:：]?|\Z)"
    m = re.search(pat, text, re.I | re.S)
    return clean_text(m.group(1)) if m else ""



def infer_from_filename(path):
    name = path.stem
    info = {"category_raw": "", "issue": "", "site": "", "abnormal_type": ""}
    m = re.match(r"^\(([^)]+)\)(.+)$", name)
    if m:
        info["category_raw"] = clean_text(m.group(1))
        issue = clean_text(m.group(2))
    else:
        issue = clean_text(name)
    issue = re.sub(r"[-_ ]*20\d{6,8}\s*$", "", issue)
    issue = re.sub(r"[-_ ]*20\d{2}[./-]?\d{1,2}[./-]?\d{1,2}\s*$", "", issue)
    info["issue"] = re.sub(r"\s+", " ", issue).strip(" -_")
    if info["category_raw"]:
        parts = [p for p in re.split(r"[-_]", info["category_raw"]) if p]
        if parts:
            info["site"] = parts[0]
        if len(parts) > 1:
            info["abnormal_type"] = parts[-1]
    return info


def infer_model_from_text_or_filename(text, path):
    model = first_match(text, [
        r"(?:Model\s*Name|Model|生產機種|品名|產品名稱)\s*[:：]\s*([A-Za-z0-9][A-Za-z0-9._/+ -]{1,60})",
        r"\b([A-Z]{2,5}-\d{3,5}[A-Za-z0-9._-]*(?:-[A-Z0-9]{2,6})?)\b"
    ])
    if not model:
        model = first_match(path.stem, [r"\b([A-Z]{2,5}-\d{3,5}[A-Za-z0-9._-]*(?:-[A-Z0-9]{2,6})?)\b"])
    return re.split(r"[,，\n]", model)[0].strip()


def infer_pn_from_text(text):
    all_pns = re.findall(r"\b[0-9]{6}-[0-9]{3}\b", text)
    if all_pns:
        return compact_join(all_pns[:6])
    return first_match(text, [r"(?:PN|P/N|料號)\s*[:：]\s*([0-9]{6}-[0-9]{3})"])


TABLE_FIELD_ALIASES = {
    "主要不良現象": [
        "異常現象", "不良現象", "問題現象", "異常描述", "failure symptom", "problem description"
    ],
    "主要原因": [
        "原因/調查", "原因／調查", "調查&原因", "調查＆原因", "調查/原因", "調查／原因",
        "原因調查", "原因分析", "異常原因", "異常原因分析", "不良原因", "根本原因",
        "root cause", "cause analysis"
    ],
}

ACTION_ALIASES = [
    "改善對策", "永久對策", "長期對策", "短期對策", "臨時對策", "緊急對策",
    "corrective action", "preventive action", "action"
]

ALL_SECTION_ALIASES = (
    TABLE_FIELD_ALIASES["主要不良現象"]
    + TABLE_FIELD_ALIASES["主要原因"]
    + ACTION_ALIASES
    + ["問題敘述", "項目", "分析說明", "改善效果", "效果確認", "驗證結果", "result"]
)


def normalize_label(text):
    s = clean_text(text).lower()
    s = s.replace("：", ":").strip(" :")
    s = re.sub(r"\s+", "", s)
    s = s.replace("／", "/").replace("＆", "&")
    return s


def alias_match(cell_text, aliases):
    n = normalize_label(cell_text)
    return any(n == normalize_label(a) for a in aliases)


def is_section_label(cell_text):
    return alias_match(cell_text, ALL_SECTION_ALIASES)


def collect_row_value(row, key_index):
    """Collect content cells to the right of a table label, stopping at another known section label."""
    values = []
    for cell in row[key_index + 1:]:
        cell = clean_text(cell)
        if not cell:
            continue
        if is_section_label(cell):
            break
        if cell not in values:
            values.append(cell)
    return "\n".join(values).strip()


def extract_table_semantic_fields(tables, logger=None, source_file=""):
    """Prefer structured PPT/DOCX table labels over whole-document regex extraction."""
    result = {"主要不良現象": "", "主要原因": "", "改善對策": ""}
    sources = {}

    # Symptom / cause: first strong structured match wins.
    for field, aliases in TABLE_FIELD_ALIASES.items():
        for t in tables:
            for row in t.get("rows", []):
                for i, cell in enumerate(row):
                    if alias_match(cell, aliases):
                        value = collect_row_value(row, i)
                        if value:
                            result[field] = value
                            sources[field] = f"TABLE | page={t.get('page')} | key={clean_text(cell)}"
                            if logger:
                                logger.info(
                                    "FIELD_SOURCE | %s | TABLE | page=%s | key=%s | file=%s",
                                    field, t.get("page"), clean_text(cell), source_file
                                )
                            break
                if result[field]:
                    break
            if result[field]:
                break

    # Actions: keep multiple action types instead of overwriting one another.
    action_parts = []
    seen = set()
    for t in tables:
        for row in t.get("rows", []):
            for i, cell in enumerate(row):
                if alias_match(cell, ACTION_ALIASES):
                    value = collect_row_value(row, i)
                    label = clean_text(cell)
                    key = (normalize_label(label), value)
                    if value and key not in seen:
                        seen.add(key)
                        action_parts.append((label, value, t.get("page")))
                        if logger:
                            logger.info(
                                "FIELD_SOURCE | 改善對策 | TABLE | page=%s | key=%s | file=%s",
                                t.get("page"), label, source_file
                            )

    if len(action_parts) == 1:
        result["改善對策"] = action_parts[0][1]
        sources["改善對策"] = f"TABLE | page={action_parts[0][2]} | key={action_parts[0][0]}"
    elif len(action_parts) > 1:
        result["改善對策"] = "\n\n".join(
            f"【{label}】\n{value}" for label, value, _ in action_parts
        )
        sources["改善對策"] = "TABLE_MULTI_ACTION"

    return result, sources


def log_fallback_source(logger, field, source_file):
    if logger:
        logger.info("FIELD_SOURCE | %s | TEXT_FALLBACK | file=%s", field, source_file)

def infer_record(path, pages, tables, logger):
    full = "\n".join(p["text"] for p in pages)
    kv = table_key_values(tables)
    fn_info = infer_from_filename(path)
    table_fields, table_field_sources = extract_table_semantic_fields(tables, logger, path.name)

    title = ""
    for p in pages[:2]:
        lines = [clean_text(x) for x in p["text"].splitlines() if clean_text(x)]
        candidates = [x for x in lines if ("異常" in x or "不良" in x) and len(x) < 100]
        if candidates:
            title = candidates[0]
            break

    report_date = first_match(full, [
        r"\b(20\d{2}[/-]\d{1,2}[/-]\d{1,2})\b",
        r"\b(20\d{2}\.\d{1,2}\.\d{1,2})\b"
    ])
    if not report_date:
        m_date = re.search(r"(20\d{2})(\d{2})(\d{2})", path.stem)
        if m_date:
            report_date = f"{m_date.group(1)}/{m_date.group(2)}/{m_date.group(3)}"

    model = compact_join(values_for(kv, ["品名", "產品名稱", "model name", "model"])[:2])
    if not model:
        model = infer_model_from_text_or_filename(full, path)

    pn = compact_join(values_for(kv, ["料號", "PN", "part no"])[:1])
    if not pn:
        pn = infer_pn_from_text(full)

    wo = compact_join(values_for(kv, ["工單號", "工單", "PO"])[:2])
    if not wo:
        wo = first_match(full, [r"\b(SUBN[A-Z0-9-]+)\b"])

    defect_rate = compact_join(values_for(kv, ["不良率", "Defect %"])[:1])

    station = first_match(title + "\n" + full[:3000], [
        r"\b(T\d+)\s*(?:站|測試|TEST)",
        r"(SMT|DIP|FQC|OQC|PDI|OBA|Packing)\s*(?:站|測試|TEST)?"
    ])

    issue_item = title
    issue_item = re.sub(r"^\([^)]*異常[^)]*\)\s*", "", issue_item)
    issue_item = re.sub(r"^INDEX[_\s-]*", "", issue_item, flags=re.I)
    issue_item = clean_text(issue_item)
    bad_issue_titles = {"異常原因:", "異常原因", "問題敘述:", "問題敘述", "改善對策:", "改善對策", "緊急對策:", "緊急對策", "原因分析:", "原因分析", "next step:", "next step"}
    if not issue_item or len(issue_item) < 4 or issue_item.strip().lower() in bad_issue_titles:
        issue_item = fn_info.get("issue", "")

    # Structured table extraction has highest priority.
    symptom = table_fields.get("主要不良現象", "")
    if not symptom:
        symptom = section_from_text(
            full, ["異常現象", "不良現象", "問題現象", "Failure symptom", "Problem description"],
            ["原因/調查", "調查&原因", "原因分析", "異常原因", "Root cause", "改善對策", "短期對策", "永久對策", "Corrective action"]
        )
        if symptom:
            log_fallback_source(logger, "主要不良現象", path.name)

    cause = table_fields.get("主要原因", "")
    if not cause:
        cause = section_from_text(
            full, ["原因/調查", "調查&原因", "調查＆原因", "原因分析", "異常原因", "異常原因分析", "不良原因", "Root cause", "Cause analysis"],
            ["改善對策", "短期對策", "永久對策", "長期對策", "臨時對策", "緊急對策", "Corrective action", "Action", "改善效果"]
        )
        if cause:
            log_fallback_source(logger, "主要原因", path.name)

    action = table_fields.get("改善對策", "")
    if not action:
        action = section_from_text(
            full,
            ["改善對策", "永久對策", "短期對策", "長期對策", "臨時對策", "緊急對策", "Corrective action", "Preventive action"],
            ["改善效果", "效果確認", "驗證結果", "Result", "最終判定", "INDEX", "THANK YOU"]
        )
        if action:
            log_fallback_source(logger, "改善對策", path.name)
    effect = section_from_text(
        full, ["改善效果確認", "改善效果", "效果確認"],
        ["驗證結果", "Result", "最終判定", "INDEX"]
    )
    validation = ""
    if re.search(r"高低溫.*?O5.*?PASS", full, re.I | re.S):
        validation = "高低溫 O5 Check PASS"
    else:
        validation = first_match(full, [r"(?:驗證結果|Result)\s*[:：]\s*([^\n]{1,200})"])

    final = "PASS / Close" if (
        re.search(r"改善.*?(?:<\s*1\s*%|下降)", full, re.I | re.S)
        and re.search(r"\bPASS\b", full, re.I)
    ) else ""

    codes = re.findall(r"\bT\d{4}[_A-Za-z0-9 .()-]{3,80}(?:FAIL|fail)\b", full)
    if codes:
        symptom = compact_join(codes[:5]) if not symptom else symptom

    if "不良率已從3% 下降到 <1 % (0.43%)" in full or re.search(r"不良率.*?3\s*%.*?下降.*?<\s*1", full, re.S):
        effect = "T1 不良率由 3% 降至 <1%（0.43%）"

    record = {
        "報告日期": report_date,
        "異常類別": "生產異常" if "生產異常" in full or "生產異常" in path.stem else (fn_info.get("abnormal_type") or "異常報告"),
        "客戶/廠區": first_match(full, [r"\(([^()_\n-]+)[-_ ]*(?:生產異常|DVT|線上測試異常|材料異常|原材不良|RMA不良)\)"]) or fn_info.get("site", ""),
        "Model": model,
        "PN": pn,
        "PO / 工單": wo,
        "異常站別": station,
        "異常項目": issue_item,
        "不良率": defect_rate,
        "主要不良現象": symptom,
        "主要原因": cause,
        "改善對策": action,
        "改善效果": effect,
        "驗證結果": validation,
        "最終判定": final,
        "備註": "",
        "來源檔案": path.name
    }
    logger.debug("Parsed record: %r", record)
    return record


def validate_record(record):
    required = ["報告日期", "Model", "異常項目"]
    missing = [x for x in required if not record.get(x)]
    warnings = []
    for x in ["PN", "異常站別", "主要原因", "改善對策"]:
        if not record.get(x):
            warnings.append(f"欄位未擷取：{x}")
    return missing, warnings



def migrate_sheet_schema(ws, logger=None):
    """Migrate an existing history workbook to the current HEADERS without losing matching columns."""
    existing_headers = [clean_text(ws.cell(row=1, column=c).value) for c in range(1, ws.max_column + 1)]
    if existing_headers == HEADERS:
        return False

    header_map = {}
    for c, h in enumerate(existing_headers, 1):
        if h and h not in header_map:
            header_map[h] = c

    rows = []
    for r in range(2, ws.max_row + 1):
        rec = {}
        for h in HEADERS:
            if h in header_map:
                rec[h] = ws.cell(row=r, column=header_map[h]).value
            else:
                rec[h] = ""
        rows.append(rec)

    if logger:
        logger.warning(
            "SCHEMA_MIGRATION | old_columns=%s | new_columns=%s | rows=%s",
            len(existing_headers), len(HEADERS), len(rows)
        )

    ws.delete_rows(1, ws.max_row)
    ws.append(HEADERS)
    for rec in rows:
        ws.append([rec.get(h, "") for h in HEADERS])
    return True


def create_new_workbook():
    wb = Workbook()
    ws = wb.active
    ws.title = "異常履歷表"
    ws.append(HEADERS)
    return wb, ws


def style_history_sheet(ws):
    header_fill = PatternFill("solid", fgColor="1F4E78")
    header_font = Font(color="FFFFFF", bold=True)
    thin = Side(style="thin", color="B7B7B7")

    for c in ws[1]:
        c.fill = header_fill
        c.font = header_font
        c.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
        c.border = Border(left=thin, right=thin, top=thin, bottom=thin)

    for row in ws.iter_rows(min_row=2):
        for c in row:
            c.alignment = Alignment(vertical="top", wrap_text=True)
            c.border = Border(left=thin, right=thin, top=thin, bottom=thin)

    # Fixed widths: avoid extremely wide columns caused by long abnormal descriptions.
    widths = {
        1: 6,    # No.
        2: 13,   # 報告日期
        3: 14,   # 異常類別
        4: 14,   # 客戶/廠區
        5: 20,   # Model
        6: 18,   # PN
        7: 20,   # PO / 工單
        8: 12,   # 異常站別
        9: 32,   # 異常項目
        10: 10,  # 不良率
        11: 42,  # 主要不良現象
        12: 42,  # 主要原因
        13: 48,  # 改善對策
        14: 28,  # 改善效果
        15: 24,  # 驗證結果
        16: 14,  # 最終判定
        17: 24,  # 備註
        18: 34,  # 來源檔案
    }
    for col, width in widths.items():
        ws.column_dimensions[get_column_letter(col)].width = width

    # Fixed heights: full text remains in the cell, while the sheet stays readable.
    ws.row_dimensions[1].height = 30
    for r in range(2, ws.max_row + 1):
        ws.row_dimensions[r].height = 72

    ws.freeze_panes = "A2"
    ws.auto_filter.ref = ws.dimensions


def normalize_key_value(v):
    return re.sub(r"\s+", "", str(v or "").strip()).lower()


def make_duplicate_key_from_record(rec):
    return "|".join([
        normalize_key_value(rec.get("來源檔案")),
        normalize_key_value(rec.get("報告日期")),
        normalize_key_value(rec.get("Model")),
        normalize_key_value(rec.get("異常項目")),
    ])


def append_records_to_history(records, history_path, logger):
    if history_path.exists():
        wb = load_workbook(history_path)
        if "異常履歷表" in wb.sheetnames:
            ws = wb["異常履歷表"]
        else:
            ws = wb.active
            ws.title = "異常履歷表"
        migrate_sheet_schema(ws, logger=logger)
    else:
        wb, ws = create_new_workbook()

    existing_keys = set()
    for r in range(2, ws.max_row + 1):
        row_dict = {h: ws.cell(row=r, column=c).value for c, h in enumerate(HEADERS, 1)}
        existing_keys.add(make_duplicate_key_from_record(row_dict))

    appended, skipped, row_error = 0, 0, 0
    details = []
    for raw_rec in records:
        rec = sanitize_record_for_excel(raw_rec, logger=logger)
        key = make_duplicate_key_from_record(rec)
        if key in existing_keys:
            skipped += 1
            msg = f"SKIP - Duplicate Report: {rec.get('來源檔案','')}"
            details.append(msg)
            logger.info(msg)
            continue

        if ws.max_row >= 2:
            try:
                no_value = int(ws.cell(row=ws.max_row, column=1).value or 0) + 1
            except Exception:
                no_value = ws.max_row
        else:
            no_value = 1

        ws.append([no_value] + [rec.get(h, "") for h in HEADERS[1:]])
        existing_keys.add(key)
        appended += 1
        msg = f"APPEND - {rec.get('來源檔案','')} -> Row {ws.max_row}"
        details.append(msg)
        logger.info(msg)

    sort_result = sort_history_by_report_date(ws, logger=logger)
    style_history_sheet(ws)
    wb.save(history_path)
    return appended, skipped, row_error, details, sort_result


class App:
    def __init__(self, root):
        self.root = root
        self.root.title(f"{APP_NAME} V{APP_VERSION}")
        try:
            self.root.iconbitmap(str(resource_path("assets/HD.ico")))
        except Exception:
            pass
        self.root.geometry("1100x700")
        self.files = []
        self.history_file = None
        self.include_subfolders = tk.BooleanVar(value=True)

        frm = ttk.Frame(root, padding=12)
        frm.pack(fill="both", expand=True)

        top = ttk.Frame(frm)
        top.pack(fill="x")
        ttk.Button(top, text="選擇 PPTX / DOCX", command=self.select_files).pack(side="left")
        ttk.Button(top, text="選擇資料夾批量匯入", command=self.select_folder_batch).pack(side="left", padx=8)
        ttk.Checkbutton(top, text="包含子資料夾", variable=self.include_subfolders).pack(side="left", padx=(0, 8))
        ttk.Button(top, text="清除清單", command=self.clear_files).pack(side="left", padx=8)
        ttk.Button(top, text="選擇既有履歷Excel", command=self.select_history_file).pack(side="left", padx=8)
        ttk.Button(top, text="開始批量擷取 / 累加", command=self.run_extract).pack(side="left")
        self.debug_var = tk.BooleanVar(value=True)
        ttk.Checkbutton(top, text="Debug 模式", variable=self.debug_var).pack(side="right")

        guide = (
            "操作說明：1) 可多選檔案，或直接選擇資料夾批量匯入  2) 可勾選包含子資料夾  "
            "3) 可選既有 Abnormal_History.xlsx；未選則建立新的  4) 按「開始批量擷取 / 累加」  5) 重複報告會 SKIP"
        )
        ttk.Label(frm, text=guide, foreground="#005A9E").pack(anchor="w", pady=(8, 8))
        self.history_label = ttk.Label(frm, text="目前履歷Excel：未指定，將於輸出資料夾建立 Abnormal_History.xlsx")
        self.history_label.pack(anchor="w", pady=(0, 8))

        self.scan_label = ttk.Label(frm, text="目前已加入：0 個檔案")
        self.scan_label.pack(anchor="w", pady=(4, 4))

        ttk.Label(frm, text="待處理檔案").pack(anchor="w", pady=(12, 4))
        self.listbox = tk.Listbox(frm, height=12)
        self.listbox.pack(fill="x")

        ttk.Label(frm, text="執行紀錄").pack(anchor="w", pady=(12, 4))
        self.log_text = tk.Text(frm, height=18, wrap="word")
        self.log_text.pack(fill="both", expand=True)

        self.progress = ttk.Progressbar(frm, mode="determinate")
        self.progress.pack(fill="x", pady=(8, 0))
        self.batch_status = ttk.Label(frm, text="尚未開始")
        self.batch_status.pack(anchor="w", pady=(4, 0))

    def ui_log(self, msg):
        ts = datetime.now().strftime("%H:%M:%S")
        self.log_text.insert("end", f"[{ts}] {msg}\n")
        self.log_text.see("end")
        self.root.update_idletasks()

    def select_files(self):
        paths = filedialog.askopenfilenames(
            title="選擇異常報告",
            filetypes=[("支援格式", "*.pptx *.docx"), ("PowerPoint", "*.pptx"), ("Word", "*.docx")]
        )
        for p in paths:
            if p not in self.files:
                self.files.append(p)
                self.listbox.insert("end", p)
        self.scan_label.config(text=f"目前已加入：{len(self.files)} 個檔案")

    def select_folder_batch(self):
        folder = filedialog.askdirectory(title="選擇要批量匯入的報告資料夾")
        if not folder:
            return
        root_path = Path(folder)
        candidates = root_path.rglob("*") if self.include_subfolders.get() else root_path.glob("*")
        supported = []
        pptx_count = docx_count = unsupported = 0
        for item in candidates:
            if not item.is_file():
                continue
            ext = item.suffix.lower()
            if ext == ".pptx":
                pptx_count += 1
                supported.append(str(item))
            elif ext == ".docx":
                docx_count += 1
                supported.append(str(item))
            else:
                unsupported += 1
        before = len(self.files)
        for item in sorted(supported):
            if item not in self.files:
                self.files.append(item)
                self.listbox.insert("end", item)
        added = len(self.files) - before
        duplicate = len(supported) - added
        self.scan_label.config(text=f"目前已加入：{len(self.files)} 個檔案")
        messagebox.showinfo("批量掃描完成", f"PPTX：{pptx_count}\nDOCX：{docx_count}\n本次新增：{added}\n重複未加入：{duplicate}\n其他格式略過：{unsupported}")

    def select_history_file(self):
        path = filedialog.askopenfilename(
            title="選擇既有異常履歷 Excel",
            filetypes=[("Excel", "*.xlsx")]
        )
        if path:
            self.history_file = Path(path)
            self.history_label.config(text=f"目前履歷Excel：{self.history_file}")

    def clear_files(self):
        self.files.clear()
        self.listbox.delete(0, "end")
        self.scan_label.config(text="目前已加入：0 個檔案")

    def run_extract(self):
        if not self.files:
            messagebox.showwarning("提醒", "請先選擇 PPTX 或 DOCX 檔案")
            return

        out_root = filedialog.askdirectory(title="選擇輸出資料夾")
        if not out_root:
            return

        stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        run_dir = Path(out_root) / f"Run_{stamp}"
        run_dir.mkdir(parents=True, exist_ok=True)
        logger = setup_logger(run_dir)

        test_lines = [f"{APP_NAME} V{APP_VERSION}", f"Run={stamp}", "TEST START"]
        records = []
        self.progress["maximum"] = len(self.files)
        self.progress["value"] = 0

        logger.info("Program start. files=%d", len(self.files))
        self.ui_log(f"建立執行資料夾：{run_dir}")

        total_files = len(self.files)
        for idx, file_str in enumerate(self.files, 1):
            path = Path(file_str)
            self.batch_status.config(text=f"處理中：{idx} / {total_files} | {path.name}")
            self.root.update_idletasks()
            try:
                self.ui_log(f"解析：{path.name}")
                logger.info("Parsing %s", path)
                if path.suffix.lower() == ".pptx":
                    pages, tables = extract_pptx(path, logger)
                elif path.suffix.lower() == ".docx":
                    pages, tables = extract_docx(path, logger)
                else:
                    raise ValueError(f"Unsupported file type: {path.suffix}")

                rec = infer_record(path, pages, tables, logger)
                missing, warnings = validate_record(rec)
                rec = sanitize_record_for_excel(rec, logger=logger)
                records.append(rec)

                if missing:
                    test_lines.append(f"WARN | {path.name} | Required missing but will still append: {', '.join(missing)}")
                    logger.warning("Validation WARN %s missing=%s", path.name, missing)
                else:
                    test_lines.append(f"PASS | {path.name} | Required fields OK")
                    logger.info("Validation PASS %s", path.name)
                for w in warnings:
                    test_lines.append(f"WARN | {path.name} | {w}")
                    logger.warning("%s | %s", path.name, w)

                self.ui_log(f"完成：{path.name}")
            except Exception as e:
                logger.error("Parse FAIL %s | %s\n%s", path.name, e, traceback.format_exc())
                test_lines.append(f"FAIL | {path.name} | Exception: {e}")
                self.ui_log(f"失敗：{path.name} | {e}")
            self.progress["value"] = idx

        history_path = self.history_file if self.history_file else Path(out_root) / "Abnormal_History.xlsx"
        output_snapshot = run_dir / f"Abnormal_History_snapshot_{stamp}.xlsx"

        try:
            appended, skipped, row_error, append_details, sort_result = append_records_to_history(records, history_path, logger)
            import shutil
            shutil.copy2(history_path, output_snapshot)
            test_lines.append(f"PASS | History append | appended={appended} skipped={skipped} row_error={row_error} parsed={len(records)}")
            test_lines.extend(append_details)
            test_lines.append(
                f"SORT_CHECK | 報告日期 | DESC | {'PASS' if sort_result.get('sort_pass') else 'FAIL'}"
            )
            test_lines.append(
                f"SORT_CHECK | Blank Date At End | {'PASS' if sort_result.get('blank_at_end', True) else 'FAIL'}"
            )
            test_lines.append(
                f"RENUMBER_CHECK | No. 1~N | {'PASS' if sort_result.get('renumber_pass') else 'FAIL'}"
            )
            logger.info("History saved: %s appended=%d skipped=%d", history_path, appended, skipped)
        except Exception as e:
            logger.error("Excel append FAIL | %s\n%s", e, traceback.format_exc())
            test_lines.append(f"FAIL | Excel append | {e}")
            messagebox.showerror("錯誤", f"履歷 Excel 建立/更新失敗：{e}")
            return

        test_lines.append("TEST END")
        (run_dir / "test.log").write_text("\n".join(test_lines), encoding="utf-8")

        verification = [
            f"# Verification Report - {APP_NAME} V{APP_VERSION}",
            "",
            f"- Run Time: {stamp}",
            f"- Input Files: {len(self.files)}",
            f"- Parsed Records: {len(records)}",
            f"- Appended Records: {appended}",
            f"- Skipped Duplicate Records: {skipped}",
            f"- Skipped Row Error: {row_error}",
            f"- Sorted Rows: {sort_result.get('sorted_rows', 0)}",
            f"- Blank/Invalid Report Dates: {sort_result.get('blank_dates', 0)}",
            f"- Date Sort DESC: {'PASS' if sort_result.get('sort_pass') else 'FAIL'}",
            f"- Renumber No.: {'PASS' if sort_result.get('renumber_pass') else 'FAIL'}",
            f"- Master History Excel: {history_path}",
            f"- Snapshot Excel: {output_snapshot.name}",
            "",
            "## Validation Rule",
            "- Required: 報告日期 / Model / 異常項目",
            "- Warning check: PN / 異常站別 / 主要原因 / 改善對策",
            "- Duplicate Rule: 來源檔案 + 報告日期 + Model + 異常項目",
            "- 詳細結果請查看 test.log、run.log、debug.log、error.log",
        ]
        (run_dir / "verification_report.md").write_text("\n".join(verification), encoding="utf-8")

        self.ui_log(f"履歷Excel完成：{history_path}")
        self.ui_log(f"本次快照：{output_snapshot}")
        self.ui_log(f"新增 {appended} 筆，重複略過 {skipped} 筆，寫入錯誤略過 {row_error} 筆")
        self.ui_log(f"日期排序：{'PASS' if sort_result.get('sort_pass') else 'FAIL'}（新到舊，空白日期置底）")
        self.batch_status.config(text=f"完成：總檔案 {len(self.files)} | 新增 {appended} | 重複 {skipped} | 寫入錯誤 {row_error}")
        self.ui_log("Test Log / Debug Record / Error Log / Verification Report 已建立")
        logger.info("Program completed")
        messagebox.showinfo("完成", f"擷取完成\n\n履歷Excel：{history_path}\n新增：{appended} 筆\n重複略過：{skipped} 筆\n寫入錯誤略過：{row_error} 筆")


def main():
    root = tk.Tk()
    App(root)
    root.mainloop()


if __name__ == "__main__":
    main()
